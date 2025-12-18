import os
import re
from pathlib import Path

ROOT = Path(r"c:\Users\ethan\Downloads\biotestprep")
OUT = ROOT / "exam_study_pack" / "_extracted_sources"

SOURCE_FILES = [
    ROOT / "week13" / "14_prok_prot.pptx",
    ROOT / "week13" / "Ch_16_Prok_Euk.pdf",
    ROOT / "week15" / "15_plants_fungi.pptx",
    ROOT / "week15" / "Ch_17_Plantae & Fungi.pdf",
    ROOT / "week16" / "16_animalia.pptx.pdf",
    ROOT / "week16" / "Ch_19_Animalia.pdf",
    ROOT / "week17" / "24_ecology.pptx",
    ROOT / "week17" / "Ch_36_37_Ecology.pdf",
]


def sanitize_filename(name: str) -> str:
    name = re.sub(r"[\\/:*?\"<>|]", "-", name)
    name = re.sub(r"\s+", " ", name).strip()
    return name


def extract_pdf_text(pdf_path: Path) -> str:
    try:
        import pdfplumber  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "pdfplumber is required to extract PDF text. Install with: pip install pdfplumber"
        ) from e

    chunks: list[str] = []
    with pdfplumber.open(str(pdf_path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                chunks.append(f"\n\n--- PAGE {i} ---\n{text}")
    return "\n".join(chunks).strip() + "\n"


def extract_pptx_text(pptx_path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "python-pptx is required to extract PPTX text. Install with: pip install python-pptx"
        ) from e

    prs = Presentation(str(pptx_path))
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                slide_text.append(text)
        if slide_text:
            chunks.append(f"\n\n--- SLIDE {idx} ---\n" + "\n".join(slide_text))
    return "\n".join(chunks).strip() + "\n"


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)

    missing = [p for p in SOURCE_FILES if not p.exists()]
    if missing:
        print("Missing files (skipping):")
        for p in missing:
            print(" -", p)

    for src in SOURCE_FILES:
        if not src.exists():
            continue

        out_name = sanitize_filename(src.name) + ".txt"
        out_path = OUT / out_name

        print(f"Extracting: {src} -> {out_path}")
        try:
            if src.suffix.lower() == ".pdf" and src.name.lower().endswith(".pptx.pdf"):
                # treat as PDF
                text = extract_pdf_text(src)
            elif src.suffix.lower() == ".pdf":
                text = extract_pdf_text(src)
            elif src.suffix.lower() == ".pptx":
                text = extract_pptx_text(src)
            else:
                print("  Unsupported type, skipping")
                continue

            out_path.write_text(text, encoding="utf-8")
        except Exception as e:
            print(f"  ERROR extracting {src.name}: {e}")

    print("Done")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
